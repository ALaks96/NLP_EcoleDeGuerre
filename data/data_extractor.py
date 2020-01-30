import os
import re
import zipfile
from io import StringIO
import PyPDF2
import lxml.etree
import pandas as pd
import slate3k as slate
import textract
from hachoir.metadata import extractMetadata
from hachoir.parser import createParser
from langdetect import detect
from pdf2image import convert_from_path
from pdfminer.converter import TextConverter
from pdfminer.layout import LAParams
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.pdfpage import PDFPage
from pptx import Presentation
from pytesseract import image_to_string
#from embedders.word2vec import vectorizer
from preprocessing.pre_processing import fix_text

try:
    from xml.etree.cElementTree import XML
except ImportError:
    from xml.etree.ElementTree import XML

WORD_NAMESPACE = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'
PARA = WORD_NAMESPACE + 'p'
TEXT = WORD_NAMESPACE + 't'


def pdf_extractor(path, vectors=False):
    # Open the pdf file in read binary mode.
    file_object = open(path, 'rb')

    # Create a pdf reader .
    pdf_file_reader = PyPDF2.PdfFileReader(file_object)

    try:
        creator = pdf_file_reader.getDocumentInfo()["/Author"]
    except:
        creator = "Unknown"

    current_page_number = 1
    paragraph_repo = {}
    vector = {}
    Classified = "No"

    # Reliably retrieve text from pdf
    with open(path, 'rb') as f:
        doc = slate.PDF(f)
    # Loop in all the pdf pages.
    for page in doc:
        # Get pdf page text.
        temp1 = None
        temp2 = None
        temp1 = page
        if vectors:
            temp2 = vectorizer(page, lang=detect(page))
        paragraph_repo[str(current_page_number)] = temp1
        vector[str(current_page_number)] = temp2

        # if "cid" in temp2:
        #     c = 0
        #     c = temp2.count("cid")
        #
        #     if c > 5:
        #         Classified = "Yes"

        if not paragraph_repo[str(current_page_number)]:
            # If can not extract text then use ocr lib to extract the scanned pdf file.
            try:
                paragraph_repo[str(current_page_number)] = fix_text(textract.process(path,
                                                                                     method='tesseract',
                                                                                     encoding='utf-8'))
            except TimeoutError:
                continue

        current_page_number += 1

    if vectors:
        return Classified, creator, paragraph_repo, vector
    else:
        return Classified, creator, paragraph_repo


def docx_extractor(path, vectors=False):
    document = zipfile.ZipFile(path)
    xml_content = document.read('word/document.xml')
    ## GET METADATA
    # use lxml to parse the xml file we are interested in
    try:
        doc = lxml.etree.fromstring(document.read('docProps/core.xml'))
        # retrieve creator
        ns = {'dc': 'http://purl.org/dc/elements/1.1/'}
        creator = doc.xpath('//dc:creator', namespaces=ns)[0].text
    except:
        creator = "Unknown"
    document.close()
    tree = XML(xml_content)

    doc = {}
    vector = {}
    paragraph_nb = 1
    for paragraph in tree.getiterator(PARA):
        texts = None
        text = ""
        texts = [node.text
                 for node in paragraph.getiterator(TEXT)
                 if node.text]
        if texts:
            text = ''.join(texts)
            doc[str(paragraph_nb)] = fix_text(text)
            if vectors:
                vector[str(paragraph_nb)] = vectorizer(text, lang=detect(text))
            paragraph_nb += 1

    if vectors:
        return creator, doc, vector
    else:
        return creator, doc


def ppt_extractor(path, vectors=False):
    filename = os.path.basename(path)
    paragraph_repo = {}
    vector = {}
    f = open(path, "rb")
    prs = Presentation(f)
    slide_nb = 0
    if filename.endswith(".pptx"):
        try:
            document = zipfile.ZipFile(path)
            ## GET METADATA
            # use lxml to parse the xml file we are interested in
            doc = lxml.etree.fromstring(document.read('docProps/core.xml'))
            # retrieve creator
            ns = {'dc': 'http://purl.org/dc/elements/1.1/'}
            creator = doc.xpath('//dc:creator', namespaces=ns)[0].text
        except:
            creator = "Unknown"
    else:
        creator = "Unknown"

    for slide in prs.slides:

        slide_nb += 1
        temp_text = ''

        for shape in slide.shapes:

            if hasattr(shape, "text") and shape.text.strip():
                temp_text += shape.text

        paragraph_repo[str(slide_nb)] = fix_text(temp_text)
        if vectors:
            vector[str(slide_nb)] = vectorizer(temp_text, lang=detect(text))

    if vectors:
        return creator, paragraph_repo, vector
    else:
        return creator, paragraph_repo


def txt_extractor(path, vectors=False):
    doc = {}
    vector = {}
    paragraph_nb = 1

    with open(path) as f:
        lines = f.read()

    texts = lines.strip().split("/n/n")
    for text in texts:
        doc[str(paragraph_nb)] = fix_text(text)
        if vectors:
            vector[str(paragraph_nb)] = vectorizer(text, lang=detect(text))
        paragraph_nb += 1

    if vectors:
        return doc, vector
    else:
        return doc


def img_extractor(path):
    parser = createParser(path)
    metadata = extractMetadata(parser)

    return metadata.exportDictionary()["Metadata"]


def get_arbo(location):
    # Initialize list of directories
    paths = []

    for file in os.listdir(location):
        path = str(location) + "/" + str(file)
        paths.append(path)

    return paths


def pdf_extractor2(path, vectors=False):
    # Open the pdf file in read binary mode.
    file_object = open(path, 'rb')

    # Create a pdf reader .
    pdf_file_reader = PyPDF2.PdfFileReader(file_object)

    try:
        creator = pdf_file_reader.getDocumentInfo()["/Author"]
    except:
        creator = "Unknown"

    # Get total pdf page number.
    total_page_number = pdf_file_reader.numPages

    current_page_number = 1
    paragraph_repo = {}
    vector = {}
    Classified = "No"

    # Loop in all the pdf pages.
    while current_page_number < total_page_number:
        # Get the specified pdf page object.
        pdf_page = pdf_file_reader.getPage(current_page_number)
        text = pdf_page.extractText()

        # Get pdf page text.
        paragraph_repo[str(current_page_number)] = text
        if vectors:
            vector[str(current_page_number)] = vectorizer(text, lang=detect(text))

        # Process next page.
        current_page_number += 1

    if not paragraph_repo:
        # If can not extract text then use ocr lib to extract the scanned pdf file.
        paragraph_repo[str(current_page_number)] = textract.process(path, method='tesseract', encoding='utf-8')

    if vectors:
        return Classified, creator, paragraph_repo, vector
    else:
        return Classified, creator, paragraph_repo


def pdf_extractor3(path, vectors=False):
    rsrcmgr = PDFResourceManager()
    retstr = StringIO()
    laparams = LAParams()
    device = TextConverter(rsrcmgr, retstr, laparams=laparams)
    fp = open(path, 'rb')
    interpreter = PDFPageInterpreter(rsrcmgr, device)
    password = ""
    maxpages = 0
    caching = True
    pagenos = set()
    creator = "Unknown"

    current_page_number = 1
    paragraph_repo = {}
    vector = {}
    Classified = False

    for page in PDFPage.get_pages(fp, pagenos, maxpages=maxpages, password=password, caching=caching, check_extractable=True):
        text = ''
        interpreter.process_page(page)

        text = retstr.getvalue()
        retstr.truncate(0)
        text = re.sub(u'(\u0000)', "", text)
        paragraph_repo[str(current_page_number)] = text
        if vectors:
            vector[str(current_page_number)] = vectorizer(text, lang=detect(text))
        current_page_number += 1

    fp.close()
    device.close()
    retstr.close()
    if vectors:
        return Classified, creator, paragraph_repo, vector
    else:
        return Classified, creator, paragraph_repo


def excel_extractor(path):
    OrderedDic = pd.read_excel(path, sheet_name=None)
    Dic = dict(OrderedDic)
    for k, v in Dic.items():
        if v:
            Dic[k] = v.dropna().to_string()
        else:
            pass
    filename = os.path.basename(path)
    if filename.endswith(".xlsx"):
        try:
            document = zipfile.ZipFile(path)
            ## GET METADATA
            # use lxml to parse the xml file we are interested in
            doc = lxml.etree.fromstring(document.read('docProps/core.xml'))
            # retrieve creator
            ns = {'dc': 'http://purl.org/dc/elements/1.1/'}
            creator = doc.xpath('//dc:creator', namespaces=ns)[0].text
        except:
            creator = "Unknown"
    else:
        creator = "Unknown"

    return creator, Dic


def scan_extractor(path, vectors=False):
    pages = convert_from_path(path, 500)
    paragraph_repo = {}
    vector = {}
    current_page_number = 1

    for page in pages:
        text = ''
        text = image_to_string(page)
        paragraph_repo[str(current_page_number)] = text

        if vectors:
            vector[str(current_page_number)] = vectorizer(text, lang=detect(text))

    if vectors:
        return paragraph_repo, vector

    else:
        return paragraph_repo